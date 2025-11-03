import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
import io

st.title("Market Cap PPT Generator")

st.markdown("Upload 7 Excel model files and the PowerPoint template:")

uploaded_excel_files = st.file_uploader("Upload Excel Model Files", type=["xlsx"], accept_multiple_files=True)
uploaded_pptx_file = st.file_uploader("Upload PowerPoint File", type=["pptx"])

PRICE_KEYWORDS = ['close', 'px', 'price', 'closing price', 'last price']
SHARES_KEYWORDS = ['shares', 'shares outstanding', 'shares os', 'share count', 'total shares']

def find_market_cap(file):
    try:
        xl = pd.ExcelFile(file)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name, header=None)

            price = None
            shares = None

            for i in range(df.shape[0]):
                for j in range(df.shape[1]):
                    cell = str(df.iat[i, j]).strip().lower()

                    if any(k in cell for k in PRICE_KEYWORDS):
                        try:
                            price = float(df.iat[i, j+1])
                        except:
                            pass
                    if any(k in cell for k in SHARES_KEYWORDS):
                        try:
                            shares = float(df.iat[i, j+1])
                        except:
                            pass

                    if price and shares:
                        return price * shares
    except Exception as e:
        st.error(f"Error processing {file.name}: {e}")
    return None

def populate_table(table, market_caps):
    # First row is header, start from second row
    row_idx = 1
    for ticker, market_cap in market_caps:
        if row_idx >= len(table.rows):
            break  # No more empty rows in table

        row = table.rows[row_idx]
        row.cells[0].text = ticker
        row.cells[1].text = f"${market_cap:,.2f}"

        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

        row_idx += 1

if uploaded_excel_files and uploaded_pptx_file:
    st.success(f"{len(uploaded_excel_files)} Excel files and 1 PowerPoint template uploaded successfully!")

    market_caps = []

    for file in uploaded_excel_files:
        market_cap = find_market_cap(file)
        if market_cap:
            ticker = file.name.split(" ")[0].strip()
            market_caps.append((ticker, market_cap))
        else:
            st.warning(f"Market Cap data not found in {file.name}")

    if market_caps:
        prs = Presentation(uploaded_pptx_file)

        slide = prs.slides[0]

        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break

        if table is None:
            st.error("No table found in the PowerPoint slide.")
        else:
            populate_table(table, market_caps)

            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)

            st.success("Presentation generated successfully!")

            st.download_button(
                label="Download Updated PowerPoint",
                data=pptx_io,
                file_name="Updated_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.error("No valid market cap data found in the uploaded Excel files.")
